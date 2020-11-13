import os
import pdftotext
from pptx import Presentation

from django.shortcuts import render
from django.views.generic.base import View
from django.views.generic.edit import CreateView
from django.views.generic.list import ListView

from ayen_task.core.forms import SearchForm
from ayen_task.core.models import Document


class Home(ListView):
    model = Document
    template_name = "home.html"


class Search(View):
    def get(self, request, *args, **kwargs):
        keyword = request.GET.get("keyword")
        object_list = []

        if keyword:
            for document in Document.objects.all():
                # searching file name
                if keyword in document.name:
                    object_list.append(document)
                # searching file content
                else:
                    path = document.file.path
                    extension = os.path.splitext(path)[1]

                    if extension == ".pdf":
                        with open(path, "rb") as f:
                            pdf = pdftotext.PDF(f)
                            for page in pdf:
                                if keyword in page:
                                    object_list.append(document)

                    elif extension == ".pptx":
                        presentation = Presentation(str(path))
                        for slide in presentation.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    shape.text = shape.text.lower()
                                    if keyword.lower() in shape.text:
                                        object_list.append(document)

        context = {"form": SearchForm(), "object_list": object_list}
        return render(request, "search.html", context)


class Upload(CreateView):
    model = Document
    fields = ["name", "file"]
    success_url = "/"
    template_name = "upload.html"
